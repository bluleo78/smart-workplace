"""파일 바이트에서 평문 텍스트를 추출. mime 별 라이브러리 분기. 순수 함수(외부 IO 없음).

지원 mime 목록은 api `com.workplace.fileai.ExtractableTypes` 와 1:1 미러다 —
한쪽만 고치면 드리프트(#735). 목록 변경 시 양쪽을 함께 수정한다.
"""
import io
from html.parser import HTMLParser
from xml.etree import ElementTree
from zipfile import ZipFile

from . import soffice


class _HtmlTextStripper(HTMLParser):
    """HTML 마크업을 제거하고 텍스트 노드만 수집. script/style 내부 텍스트는 버린다.

    표준 라이브러리(html.parser)만 사용 — 신규 의존성 없음.
    """

    def __init__(self):
        super().__init__()
        self._chunks: list[str] = []
        self._skip_depth = 0  # script/style 내부 깊이(중첩 대비 카운터)

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip_depth += 1

    def handle_endtag(self, tag):
        if tag in ("script", "style") and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data):
        if self._skip_depth == 0:
            self._chunks.append(data)

    def get_text(self) -> str:
        return " ".join(chunk.strip() for chunk in self._chunks if chunk.strip())


def _strip_html(data: bytes) -> str:
    """HTML 바이트 → 마크업 제거된 평문. 원문 그대로 넘기면 요약 토큰을 마크업이 잠식한다."""
    stripper = _HtmlTextStripper()
    stripper.feed(data.decode("utf-8", errors="replace"))
    return stripper.get_text().strip()


def _extract_pptx(data: bytes) -> str:
    """PPTX: python-pptx 로 슬라이드 텍스트 + 노트를 개행 결합."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text:
                    parts.append(text)
        # 노트가 없는 슬라이드도 있으므로 존재할 때만 수집.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                parts.append(notes)
    return "\n".join(parts).strip()


def _extract_hwpx(data: bytes) -> str:
    """HWPX = ZIP + XML. Contents/section*.xml 의 모든 텍스트 노드를 수집.

    태그에 네임스페이스가 붙으므로(hp:t 등) 태그명으로 필터하지 않고
    표준 라이브러리(zipfile + xml.etree)만으로 전체 텍스트 노드를 순회한다.
    """
    parts: list[str] = []
    with ZipFile(io.BytesIO(data)) as zf:
        section_names = sorted(
            n for n in zf.namelist() if n.startswith("Contents/section") and n.endswith(".xml")
        )
        for name in section_names:
            root = ElementTree.fromstring(zf.read(name))
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    parts.append(elem.text.strip())
    return "\n".join(parts).strip()


# 레거시 OLE·한글(.hwp) → soffice 변환 시 사용할 확장자 매핑.
_SOFFICE_SUFFIX_BY_MIME = {
    "application/msword": ".doc",
    "application/vnd.ms-powerpoint": ".ppt",
    "application/vnd.ms-excel": ".xls",
    "application/x-hwp": ".hwp",
}

# 텍스트 계열 범용 처리에 포함되는 application/* mime(Global Constraints 표와 동일).
_TEXT_LIKE_APPLICATION_MIMES = (
    "application/json",
    "application/xml",
    "application/x-yaml",
    "application/yaml",
    "application/javascript",
    "application/x-sh",
)


def extract_text(data: bytes, mime: str, max_chars: int) -> dict:
    """
    파일 바이트를 받아 평문 텍스트를 추출한다.

    :param data: 파일 원본 바이트
    :param mime: MIME 타입 문자열
    :param max_chars: 최대 반환 문자 수 (초과 시 잘라냄)
    :return: {text, charCount, lang, truncated}
             - text: 추출된 텍스트 (미지원 mime → "")
             - charCount: 반환된 text 길이 (I1: api ExtractResult.charCount 와 키 통일)
             - lang: 언어 코드(현재 미구현, None)
             - truncated: max_chars 초과로 잘렸으면 True
    """
    text = _dispatch(data, mime)
    truncated = len(text) > max_chars
    if truncated:
        text = text[:max_chars]
    # I1 수정: 워커→api 콜백 페이로드 키를 charCount(camelCase)로 통일.
    # 이전 char_count(snake_case)는 Jackson 기본 설정으로 역직렬화 불가(null 폴백).
    return {"text": text, "charCount": len(text), "lang": None, "truncated": truncated}


def _dispatch(data: bytes, mime: str) -> str:
    """mime 타입에 따라 적절한 파서를 선택해 텍스트 추출.

    api `ExtractableTypes` 와 1:1 미러 — 한쪽만 고치면 드리프트(#735).
    """
    if mime == "application/pdf":
        # PDF: pymupdf(fitz) 로 페이지별 텍스트 결합
        import fitz  # pymupdf
        with fitz.open(stream=data, filetype="pdf") as doc:
            return "\n".join(page.get_text() for page in doc).strip()

    if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        # DOCX: python-docx 로 단락 텍스트 추출
        from docx import Document
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs).strip()

    if mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        # XLSX: openpyxl 로 모든 시트 셀 값을 CSV 형태로 결합
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                rows.append(",".join("" if c is None else str(c) for c in row))
        return "\n".join(rows).strip()

    if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        # PPTX: python-pptx 로 슬라이드 텍스트 + 노트 추출
        return _extract_pptx(data)

    if mime == "application/hwp+zip":
        # HWPX: ZIP+XML 표준 라이브러리 파싱
        return _extract_hwpx(data)

    suffix = _SOFFICE_SUFFIX_BY_MIME.get(mime)
    if suffix is not None:
        # 레거시 OLE(.doc/.ppt/.xls)·.hwp: LibreOffice headless 변환
        return soffice.convert_to_text(data, suffix)

    if mime == "text/html":
        # HTML: 범용 text/ 분기보다 먼저 처리 — 마크업 제거 후 텍스트만 남긴다.
        return _strip_html(data)

    if mime.startswith("text/") or mime in _TEXT_LIKE_APPLICATION_MIMES:
        # 텍스트 계열: UTF-8 디코딩(오류 문자 대체)
        return data.decode("utf-8", errors="replace").strip()

    # 미지원 mime(이미지 등) → 빈 텍스트(호출측이 SKIPPED 판정)
    return ""
