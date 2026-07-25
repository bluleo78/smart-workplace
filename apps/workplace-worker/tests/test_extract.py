"""
extract_text 순수 함수 단위 테스트.
txt/csv 는 인라인 바이트, docx/xlsx 는 라이브러리로 생성.
"""
import pytest
from app.extract import extract_text


def test_plain_text():
    """일반 텍스트 추출 + charCount 검증(I1: snake→camelCase 통일)."""
    r = extract_text(b"hello world", "text/plain", max_chars=1000)
    assert r["text"] == "hello world"
    assert r["charCount"] == 11  # I1: char_count → charCount
    assert r["truncated"] is False


def test_csv_text():
    """CSV는 text/ prefix로 일반 텍스트 처리."""
    r = extract_text(b"a,b\n1,2", "text/csv", max_chars=1000)
    assert "a,b" in r["text"]


def test_truncation():
    """max_chars 초과 시 잘라내고 truncated=True."""
    r = extract_text(b"x" * 5000, "text/plain", max_chars=100)
    assert r["charCount"] == 100  # I1: char_count → charCount
    assert r["truncated"] is True


def test_unsupported_mime_returns_empty():
    """미지원 mime → 빈 텍스트 반환(호출측이 SKIPPED 판정)."""
    r = extract_text(b"\x00\x01", "application/octet-stream", max_chars=1000)
    assert r["text"] == ""


def test_unsupported_mime_zip_returns_empty():
    """application/zip 도 미지원 목록 밖 → 빈 텍스트."""
    r = extract_text(b"\x00\x01", "application/zip", max_chars=10_000)
    assert r["text"] == ""


def test_html_strips_tags_and_script():
    """HTML 은 마크업을 제거한 텍스트만 남긴다 — 원문 그대로면 요약 토큰을 마크업이 잠식한다."""
    html = b"<html><head><style>p{color:red}</style><script>alert(1)</script></head><body><h1>\xec\xa0\x9c\xeb\xaa\xa9</h1><p>\xeb\xb3\xb8\xeb\xac\xb8</p></body></html>"
    out = extract_text(html, "text/html", 10_000)
    assert "제목" in out["text"]
    assert "본문" in out["text"]
    assert "alert(1)" not in out["text"]
    assert "color:red" not in out["text"]
    assert "<h1>" not in out["text"]


def test_hwpx_extracts_section_text():
    """hwpx = ZIP + XML — 표준 라이브러리만으로 section*.xml 의 텍스트 노드를 수집한다."""
    import zipfile, io
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr(
            "Contents/section0.xml",
            '<?xml version="1.0"?><root xmlns:hp="x"><hp:t>한글 문서 본문</hp:t></root>',
        )
    out = extract_text(buf.getvalue(), "application/hwp+zip", 10_000)
    assert "한글 문서 본문" in out["text"]


def test_pptx_extracts_slide_text():
    """PPTX: python-pptx 로 생성한 픽스처를 왕복 검증."""
    from pptx import Presentation
    from pptx.util import Inches
    import io

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "발표 자료 핵심"
    buf = io.BytesIO()
    prs.save(buf)

    out = extract_text(
        buf.getvalue(),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        10_000,
    )
    assert "발표 자료 핵심" in out["text"]


def test_docx(tmp_path):
    """docx 파일에서 본문 단락 추출."""
    from docx import Document
    p = tmp_path / "d.docx"
    doc = Document()
    doc.add_paragraph("도큐먼트 본문")
    doc.save(p)
    r = extract_text(
        p.read_bytes(),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        max_chars=1000,
    )
    assert "도큐먼트 본문" in r["text"]


def test_xlsx(tmp_path):
    """xlsx 파일에서 셀 값 추출."""
    from openpyxl import Workbook
    p = tmp_path / "s.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.append(["이름", "점수"])
    ws.append(["홍길동", 100])
    wb.save(p)
    r = extract_text(
        p.read_bytes(),
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        max_chars=1000,
    )
    assert "이름" in r["text"]
    assert "홍길동" in r["text"]
